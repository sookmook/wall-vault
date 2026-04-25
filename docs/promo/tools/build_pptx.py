"""Build wall-vault-pitch.pptx directly via python-pptx.

Marp + browser export is unavailable in this environment (puppeteer
chromium needs libasound2 which we cannot install without sudo). We
keep deck.md as the human-readable narrative source and use this script
as the actual PPTX compiler. Output structure mirrors deck.md.

12 slides, 16:9, marketing-tone Korean copy, uses the brand palette
inherited from logo (navy ink + gold primary on cream paper).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


REPO = Path(__file__).resolve().parents[3]
PROMO = REPO / "docs/promo"
IMG = PROMO / "images"
OUTPUT = PROMO / "output"

# ── Brand palette (fallback to hand-tuned values if palette.css drift) ──
INK     = RGBColor(0x10, 0x22, 0x36)
MUTED   = RGBColor(0x5b, 0x68, 0x77)
PRIMARY = RGBColor(0xc8, 0x88, 0x1a)
ACCENT  = RGBColor(0xe0, 0x9a, 0x32)
PAPER   = RGBColor(0xff, 0xfa, 0xf2)
LINE    = RGBColor(0xea, 0xd8, 0xb6)
SUCCESS = RGBColor(0x2f, 0x9d, 0x6b)
DANGER  = RGBColor(0xc6, 0x4a, 0x3c)

FONT = "Noto Sans CJK KR"   # widely available CJK font on Linux; will be embedded by reference
FONT_FALLBACK = "맑은 고딕"  # Windows fallback; PowerPoint substitutes if Noto unavailable


# ── Geometry (16:9 at 1920x1080 px = 13.333 x 7.5 inches) ───────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color: RGBColor):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # send to back so subsequent shapes layer on top
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(
    slide,
    left, top, width, height,
    text: str,
    *,
    size: int = 24,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return box


def add_multi(
    slide,
    left, top, width, height,
    runs: list[tuple[str, dict]],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Several runs in one paragraph for inline emphasis."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (text, attrs) in enumerate(runs):
        if i == 0 and not p.runs:
            run = p.add_run()
        else:
            run = p.add_run()
        run.text = text
        f = run.font
        f.name = attrs.get("font", FONT)
        f.size = Pt(attrs.get("size", 24))
        f.bold = attrs.get("bold", False)
        f.color.rgb = attrs.get("color", INK)
    return box


def add_bullets(
    slide,
    left, top, width, height,
    items: list[str],
    *,
    size: int = 24,
    color: RGBColor = INK,
    line_spacing: float = 1.4,
    bullet_color: RGBColor = PRIMARY,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        # bullet glyph in primary
        b = p.add_run()
        b.text = "•  "
        b.font.name = FONT
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        b.font.bold = True
        # body
        r = p.add_run()
        r.text = item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_chip(
    slide, left, top, label: str,
    *,
    solid: bool = False,
    size: int = 16,
):
    pad_x = Inches(0.18)
    pad_y = Inches(0.08)
    # measure roughly: ~0.085" per char (good enough for Korean + Latin mix)
    width = Inches(max(0.6, len(label) * 0.13)) + pad_x * 2
    height = Inches(0.42)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    chip.adjustments[0] = 0.5
    chip.line.color.rgb = PRIMARY
    chip.line.width = Pt(1.5)
    chip.fill.solid()
    if solid:
        chip.fill.fore_color.rgb = PRIMARY
        text_color = PAPER
    else:
        chip.fill.fore_color.rgb = RGBColor(0xff, 0xf4, 0xe0)
        text_color = PRIMARY
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = pad_x
    tf.margin_right = pad_x
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return chip, width


def add_chip_row(slide, left, top, labels: list[str], *, solid=False, size=16, gap_in=0.1):
    x = left
    for lab in labels:
        _, w = add_chip(slide, x, top, lab, solid=solid, size=size)
        x = x + w + Inches(gap_in)
    return x


def add_kpi(slide, left, top, num: str, lbl: str, width=Inches(2.2), height=Inches(1.6)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.06
    box.line.color.rgb = LINE
    box.line.width = Pt(2)
    box.fill.solid()
    box.fill.fore_color.rgb = PAPER
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.name = FONT
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = lbl
    r2.font.name = FONT
    r2.font.size = Pt(16)
    r2.font.color.rgb = MUTED
    return box


def add_image(slide, left, top, path: Path, *, width=None, height=None):
    if not path.exists():
        raise FileNotFoundError(path)
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return pic


def add_footer(slide, page: int, total: int):
    add_text(
        slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        "wall-vault · 투자·제휴 피칭 · 2026",
        size=11, color=MUTED,
    )
    add_text(
        slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
        f"{page} / {total}",
        size=11, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ────────────────────────────────────────────────────────────────────────
# Slide builders — one function per slide, ordered by deck position.
# ────────────────────────────────────────────────────────────────────────


def slide_01_title(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    # subtle radial overlay (approximate via diagonal gradient fill on a large
    # rounded rect — python-pptx does support gradient stops but spec is verbose;
    # we use a simple tinted layer instead)
    halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(-1.5), Inches(10), Inches(8))
    halo.line.fill.background()
    halo.fill.solid()
    halo.fill.fore_color.rgb = RGBColor(0xfb, 0xee, 0xd2)
    halo.shadow.inherit = False
    spTree = halo._element.getparent()
    spTree.remove(halo._element)
    spTree.insert(3, halo._element)

    add_image(s, Inches(5.3), Inches(0.8), IMG / "logo.png", height=Inches(2.4))

    add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(1.4),
             "wall-vault", size=80, bold=True, color=INK, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0), Inches(4.95), SLIDE_W, Inches(0.8),
             "AI 가 절대 끊기지 않는 멀티벤더 게이트웨이",
             size=28, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)

    add_text(s, Inches(0), Inches(5.7), SLIDE_W, Inches(0.6),
             "키 금고 · 지능형 라우팅 · fleet observability — 한 바이너리에",
             size=18, color=MUTED, align=PP_ALIGN.CENTER)
    return s


def slide_02_section(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_text(s, Inches(0), Inches(2.8), SLIDE_W, Inches(0.6),
             "문제 의식", size=24, color=PRIMARY,
             align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(1.2),
             "AI 가 멈추는 순간, 일도 멈춥니다",
             size=56, color=PAPER, align=PP_ALIGN.CENTER, bold=True)
    return s


def slide_03_problem(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)

    add_multi(
        s, Inches(0.7), Inches(0.6), Inches(11.9), Inches(1.0),
        [
            ("AI 가 끊기는 순간, ", {"size": 38, "color": INK, "bold": True}),
            ("비용은 두 배가 됩니다", {"size": 38, "color": PRIMARY, "bold": True}),
        ],
    )
    # divider
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.55), Inches(12.6), Inches(1.55))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_bullets(
        s, Inches(0.7), Inches(1.95), Inches(7.2), Inches(5),
        [
            "클라우드 API 키 만료·쿨다운·credit-out — 한 키 죽으면 fleet 정지",
            "벤더별 모델 ID 네임스페이스가 달라 이식 비용 누적",
            "여러 머신·여러 에이전트 동시 운영 시 신호 가시성 0",
            "Fallback 이 조용히 다른 모델로 바꿔치기 → 응답 품질 붕괴",
            "운영자가 매일 \"왜 또 멈췄지?\" 를 묻게 됨",
        ],
        size=22, line_spacing=1.55,
    )

    # quote box
    quote = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(2.1), Inches(4.4), Inches(3.8))
    quote.line.fill.background()
    quote.fill.solid()
    quote.fill.fore_color.rgb = RGBColor(0xff, 0xf4, 0xe0)
    quote.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(2.1), Inches(0.12), Inches(3.8))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.shadow.inherit = False

    add_text(
        s, Inches(8.7), Inches(2.4), Inches(4.0), Inches(3.2),
        "“키 하나 죽으면 fleet 전체가 멈춥니다.\n그게 어디서 어떻게 멈췄는지조차\n보이지 않습니다.”",
        size=22, color=INK, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(s, page, total)
    return s


def slide_04_definition(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_text(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.7),
             "wall-vault — 한 줄 정의",
             size=28, color=MUTED, bold=True)

    add_multi(
        s, Inches(0), Inches(2.2), SLIDE_W, Inches(1.2),
        [
            ("AES-GCM 암호화 키 금고", {"size": 50, "bold": True, "color": PRIMARY}),
            ("  +  ", {"size": 50, "color": MUTED}),
            ("지능형 멀티벤더 프록시", {"size": 50, "bold": True, "color": PRIMARY}),
        ],
        align=PP_ALIGN.CENTER,
    )

    chips_top = Inches(4.0)
    labels = ["암호화 저장", "자동 키 로테이션", "정직한 라우팅", "실시간 신호등"]
    # layout chips evenly
    chip_w_est = sum(max(0.6, len(lab) * 0.13) + 0.36 for lab in labels) + 0.1 * (len(labels) - 1)
    start_x = Inches((13.333 - chip_w_est) / 2)
    add_chip_row(s, start_x, chips_top, labels, solid=True, size=20, gap_in=0.2)

    add_text(
        s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.0),
        "한 개의 Go 바이너리, 무설치 의존성, 17 개 언어, Linux / macOS / Windows / WSL.",
        size=22, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_footer(s, page, total)
    return s


def slide_05_architecture(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_text(s, Inches(0.7), Inches(0.5), Inches(11), Inches(0.7),
             "아키텍처 — 한 장으로 보는 동작",
             size=32, bold=True, color=INK)
    add_image(s, Inches(0.7), Inches(1.5), IMG / "architecture.png", height=Inches(4.6))
    add_text(
        s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.7),
        "벽금고(:56243) 가 키·설정의 단일 소스, 프록시(:56244) 가 지능형 라우터. SSE 로 실시간 동기화.",
        size=18, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_footer(s, page, total)
    return s


def slide_06_section(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_text(s, Inches(0), Inches(2.8), SLIDE_W, Inches(0.6),
             "차별점", size=24, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(1.2),
             "시중 솔루션과 다른 세 가지",
             size=56, color=PAPER, align=PP_ALIGN.CENTER, bold=True)
    return s


def slide_07_diff_multivendor(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_multi(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        [
            ("차별점 ① ", {"size": 28, "color": MUTED, "bold": True}),
            ("멀티벤더를 진짜로 통합", {"size": 32, "color": INK, "bold": True}),
        ],
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_bullets(
        s, Inches(0.7), Inches(1.85), Inches(7.5), Inches(5),
        [
            "4 가지 API 포맷 동시 노출 — Gemini · OpenAI · Anthropic · OpenRouter",
            "170+ 모델 자동 라우팅 — 로컬 Ollama 부터 Claude / Gemini / GPT 까지",
            "클라이언트는 자기가 가장 익숙한 형식 그대로 호출",
            "Cline · Cursor · Claude Code · OpenClaw — 코드 변경 없이 통합",
        ],
        size=22, line_spacing=1.55,
    )

    add_kpi(s, Inches(8.7), Inches(2.0), "4",   "API 포맷")
    add_kpi(s, Inches(11.0), Inches(2.0), "170+", "모델")
    add_kpi(s, Inches(8.7), Inches(3.9), "10+", "에이전트 타입")
    add_kpi(s, Inches(11.0), Inches(3.9), "17",  "로케일")

    add_footer(s, page, total)
    return s


def slide_08_diff_routing(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_multi(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        [
            ("차별점 ② ", {"size": 28, "color": MUTED, "bold": True}),
            ("정직한 라우팅 (Strict-by-default)", {"size": 32, "color": INK, "bold": True}),
        ],
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_image(s, Inches(0.6), Inches(1.7), IMG / "dispatch.png", height=Inches(5.0))

    add_bullets(
        s, Inches(7.5), Inches(2.0), Inches(5.5), Inches(4.5),
        [
            "Primary 실패 시 무단 모델 치환 없음",
            "Fallback 은 명시적 opt-in (FallbackServices) 만",
            "응답 헤더로 실제 사용 서비스·모델·fallback 사유 노출",
            "X-WV-Used-Service / X-WV-Used-Model / X-WV-Fallback-Reason",
            "호출자가 한 줄 헤더로 검증 가능",
        ],
        size=20, line_spacing=1.5,
    )

    add_footer(s, page, total)
    return s


def slide_09_diff_observability(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_multi(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        [
            ("차별점 ③ ", {"size": 28, "color": MUTED, "bold": True}),
            ("Fleet observability — 한 곳에서 다 봅니다", {"size": 32, "color": INK, "bold": True}),
        ],
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_image(s, Inches(0.7), Inches(1.6), IMG / "signal-lights.png", width=Inches(11.9))
    add_footer(s, page, total)
    return s


def slide_10_proof(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_text(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        "실측 운영 — 여러 머신, 다양한 OS, 끊김 없음",
        size=32, bold=True, color=INK,
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_image(s, Inches(0.4), Inches(1.7), IMG / "fleet-topology.png", height=Inches(5.0))

    add_bullets(
        s, Inches(7.0), Inches(1.9), Inches(6.0), Inches(5),
        [
            "WSL/Linux · macOS/GPU · ARM SBC · 전용 도메인 노드 — 동일 바이너리",
            "호스트 매칭 + 에이전트 타입별 liveness probe — 거짓 신호 모두 제거",
            "AES-GCM 키 로테이션 + 자동 쿨다운 — credit-out 일에도 즉시 로컬 전환",
            "단일 머신부터 수십 노드 분산까지 — 같은 도구",
        ],
        size=20, line_spacing=1.55,
    )
    add_footer(s, page, total)
    return s


def slide_11_global(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_text(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        "17 개 언어 — 글로벌 준비 완료",
        size=32, bold=True, color=INK,
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_text(
        s, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.6),
        "대시보드·CLI·시스템 메시지 전체 다국어. 언어 전환은 클릭 한 번.",
        size=22, color=MUTED,
    )

    flags = "🇰🇷  🇺🇸  🇯🇵  🇨🇳  🇩🇪  🇫🇷  🇪🇸  🇵🇹  🇮🇩  🇹🇭  🇮🇳  🇳🇵  🇲🇳  🇸🇦  🇪🇹  🇰🇪  🇿🇦"
    add_text(
        s, Inches(0.7), Inches(2.7), Inches(11.9), Inches(1.0),
        flags, size=42, align=PP_ALIGN.CENTER,
    )

    langs = ["한국어", "English", "日本語", "中文", "Deutsch", "Français", "Español",
             "Português", "Bahasa", "ไทย", "हिन्दी", "नेपाली", "Монгол",
             "العربية", "Hausa", "Kiswahili", "isiZulu"]
    # two rows of chips, centered roughly
    row1 = langs[:9]
    row2 = langs[9:]
    # use fixed gaps; eyeballed start positions
    add_chip_row(s, Inches(0.5), Inches(4.3), row1, size=15, gap_in=0.15)
    add_chip_row(s, Inches(1.5), Inches(4.95), row2, size=15, gap_in=0.15)

    add_text(
        s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.8),
        "한국 본사 + 동남아 · 중동 · 아프리카 · 라틴 시장 동시 대응 — 로컬라이즈된 AI 게이트웨이.",
        size=20, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_footer(s, page, total)
    return s


def slide_12_compare(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_text(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        "경쟁 비교 — 암호화 · 페르소나 · 실시간성 모두",
        size=30, bold=True, color=INK,
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    rows = [
        ("기능", "wall-vault", "LiteLLM", "OpenRouter SDK", "벤더 SDK 직접"),
        ("AES-GCM 암호화 키 금고",       "✓", "✗", "✗", "✗"),
        ("4 API 포맷 동시 노출",          "✓", "✓", "✗", "✗"),
        ("Strict-by-default + fallback 헤더", "✓", "✗", "✗", "✗"),
        ("SSE 실시간 fleet 동기화",       "✓", "✗", "✗", "✗"),
        ("에이전트 페르소나 · 음성 통합", "✓", "✗", "✗", "✗"),
        ("17 개국 다국어 UI",             "✓", "✗", "✗", "✗"),
        ("한 바이너리 · 무의존",          "✓", "Python", "JS", "언어별"),
    ]

    cols = 5
    rows_n = len(rows)
    table = s.shapes.add_table(rows_n, cols, Inches(0.7), Inches(1.7),
                               Inches(11.9), Inches(0.5) * rows_n).table
    # column widths
    table.columns[0].width = Inches(4.4)
    for i in range(1, 5):
        table.columns[i].width = Inches(1.875)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""  # clear default
            tf = cell.text_frame
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.name = FONT
            if r == 0:
                run.font.size = Pt(15)
                run.font.bold = True
                run.font.color.rgb = PAPER
                cell.fill.solid()
                cell.fill.fore_color.rgb = INK
            else:
                run.font.size = Pt(15)
                if c == 1 and val == "✓":
                    run.font.bold = True
                    run.font.color.rgb = SUCCESS
                    run.font.size = Pt(20)
                elif val == "✓":
                    run.font.color.rgb = SUCCESS
                    run.font.bold = True
                elif val == "✗":
                    run.font.color.rgb = DANGER
                    run.font.bold = True
                else:
                    run.font.color.rgb = INK
                if c == 0:
                    run.font.color.rgb = INK
                    run.font.bold = True
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xff, 0xf4, 0xe0)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PAPER
    add_footer(s, page, total)
    return s


def slide_13_roadmap(prs, total, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    add_text(
        s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.8),
        "로드맵 — 음성 → 자율 기동 → 멀티테넌시",
        size=32, bold=True, color=INK,
    )
    div = s.shapes.add_connector(1, Inches(0.7), Inches(1.4), Inches(12.6), Inches(1.4))
    div.line.color.rgb = PRIMARY
    div.line.width = Pt(2.5)

    add_image(s, Inches(0.7), Inches(1.6), IMG / "roadmap.png", height=Inches(5.0))
    add_footer(s, page, total)
    return s


def slide_14_cta(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)
    halo = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(-1), Inches(10), Inches(7))
    halo.line.fill.background()
    halo.fill.solid()
    halo.fill.fore_color.rgb = RGBColor(0xfb, 0xee, 0xd2)
    halo.shadow.inherit = False
    spTree = halo._element.getparent()
    spTree.remove(halo._element)
    spTree.insert(3, halo._element)

    add_text(
        s, Inches(0), Inches(2.0), SLIDE_W, Inches(1.4),
        "함께 만들어 갈까요?",
        size=72, bold=True, color=INK, align=PP_ALIGN.CENTER,
    )

    chips_top = Inches(3.8)
    labels = ["데모 요청", "파일럿 도입", "기술 제휴", "투자 검토"]
    chip_w_est = sum(max(0.6, len(lab) * 0.16) + 0.36 for lab in labels) + 0.2 * (len(labels) - 1)
    start_x = Inches((13.333 - chip_w_est) / 2)
    add_chip_row(s, start_x, chips_top, labels, solid=True, size=20, gap_in=0.25)

    add_text(
        s, Inches(0), Inches(5.4), SLIDE_W, Inches(0.8),
        "연구실 운영자에서 출발했습니다. 이제 시장 차례입니다.",
        size=22, color=MUTED, align=PP_ALIGN.CENTER,
    )
    return s


# ────────────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_01_title,
        slide_02_section,
        slide_03_problem,
        slide_04_definition,
        slide_05_architecture,
        slide_06_section,
        slide_07_diff_multivendor,
        slide_08_diff_routing,
        slide_09_diff_observability,
        slide_10_proof,
        slide_11_global,
        slide_12_compare,
        slide_13_roadmap,
        slide_14_cta,
    ]
    total = len(slides)

    for idx, builder in enumerate(slides, start=1):
        if idx in (1, 14):  # title / CTA — no footer
            builder(prs, total) if idx == 1 else builder(prs, total)
        elif idx in (2, 6):  # section breaks — no footer
            builder(prs, total, idx)
        else:
            builder(prs, total, idx)

    OUTPUT.mkdir(exist_ok=True)
    out = OUTPUT / "wall-vault-pitch.pptx"
    prs.save(str(out))
    print(f"✓ {out}  ({out.stat().st_size:,} bytes, {total} slides)")


if __name__ == "__main__":
    main()
