"""Wrap pre-rendered slide PNGs into a single .pptx.

Each slide is laid down as a full-bleed background image. The PPTX is
no longer text-editable, but visual fidelity matches the PIL renderer
exactly — no font substitution, no auto-layout reflow.

After the PPTX is written the intermediate PNG folder is deleted —
the .pptx already embeds every image and the loose PNGs are pure
build artifacts. Pass `--keep-slides` to skip cleanup (useful when
debugging a slide that looks wrong inside the .pptx).
"""

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


REPO = Path(__file__).resolve().parents[3]
SLIDES = REPO / "docs/promo/output/slides"
OUT = REPO / "docs/promo/output/wall-vault-pitch.pptx"


def main():
    keep_slides = "--keep-slides" in sys.argv

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    pngs = sorted(SLIDES.glob("slide-*.png"))
    if not pngs:
        raise SystemExit(f"no slide PNGs in {SLIDES}")

    blank_layout = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png),
            left=0, top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"✓ {OUT}  ({OUT.stat().st_size:,} bytes, {len(pngs)} slides)")

    if keep_slides:
        print(f"  (slides/ retained: {SLIDES})")
    else:
        shutil.rmtree(SLIDES, ignore_errors=True)
        print(f"  (slides/ cleaned — pass --keep-slides to retain)")


if __name__ == "__main__":
    main()
